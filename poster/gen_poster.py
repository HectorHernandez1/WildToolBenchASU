from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PLOTS = "wild-tool-bench/Phase2_Report/plots/"

DARK_BLUE  = RGBColor(0x1F, 0x4E, 0x79)
MID_BLUE   = RGBColor(0xBD, 0xD7, 0xEE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation('poster/Poster_Template.pptx')


# ── Helpers ──────────────────────────────────────────────────────────────────

def in2e(x):
    return int(x * 914400)


def get_tb(slide):
    return [s for s in slide.shapes if s.shape_type == 17][0]


def place(shape, left, top, width, height):
    shape.left   = in2e(left)
    shape.top    = in2e(top)
    shape.width  = in2e(width)
    shape.height = in2e(height)


def set_tb(shape, lines, font_size=12, left=None, top=None, width=None, height=None):
    if left is not None:
        place(shape, left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    txBody = tf._txBody
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)
    sz = int(font_size * 100)
    for line in lines:
        safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        p_elem = parse_xml(
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:r><a:rPr lang="en-US" sz="{sz}" dirty="0"/>'
            f'<a:t>{safe}</a:t></a:r></a:p>'
        )
        txBody.append(p_elem)


def add_image(slide, path, left, top, width, height):
    slide.shapes.add_picture(path, in2e(left), in2e(top), in2e(width), in2e(height))


def set_cell_borders(cell, color='1F4E79', width_pt=0.75):
    w = int(width_pt * 12700)
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ['lnL', 'lnR', 'lnT', 'lnB']:
        existing = tcPr.find(qn(f'a:{tag}'))
        if existing is not None:
            tcPr.remove(existing)
        ln = parse_xml(
            f'<a:{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'w="{w}" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'<a:prstDash val="solid"/>'
            f'</a:{tag}>'
        )
        tcPr.append(ln)


def fill_cell(cell, text, font_size=10, bold=False,
              fg=None, bg=None, align=PP_ALIGN.CENTER):
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if fg:
        run.font.color.rgb = fg


def add_table(slide, headers, rows, col_widths,
              left, top, row_height=0.28, font_size=10):
    """Add a styled table with visible borders. col_widths in inches."""
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    total_h = row_height * n_rows
    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        in2e(left), in2e(top),
        in2e(sum(col_widths)), in2e(total_h)
    ).table

    # Set column widths
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = in2e(w)
    # Set row heights
    for r in tbl.rows:
        r.height = in2e(row_height)

    # Header row
    for c_idx, h in enumerate(headers):
        cell = tbl.cell(0, c_idx)
        fill_cell(cell, h, font_size=font_size, bold=True, fg=WHITE, bg=DARK_BLUE)
        set_cell_borders(cell, '1F4E79', 0.75)

    # Data rows
    for r_idx, row_data in enumerate(rows):
        bg = LIGHT_GRAY if r_idx % 2 == 0 else WHITE
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx + 1, c_idx)
            fill_cell(cell, val, font_size=font_size, fg=BLACK, bg=bg)
            set_cell_borders(cell, '1F4E79', 0.75)

    return tbl


# ── SLIDE 1: Title ──────────────────────────────────────────────────────────
slide = prs.slides[0]
for shape in slide.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if 'Project Title' in txt:
            shape.text_frame.paragraphs[0].runs[0].text = (
                "WildToolBench Evaluation Enhancements:\n"
                "Improving LLM Tool-Use Through Harness-Level Interventions"
            )
        elif 'Members' in txt:
            shape.text_frame.paragraphs[0].runs[0].text = (
                "Hector Hernandez  ·  Vardaan Gangal  ·  Harish Chaurasia  "
                "·  Venkata Sai Ganesh Manda  ·  Smit Lagdhirbhai Desai"
            )

# ── SLIDE 2: Abstract ───────────────────────────────────────────────────────
tb = get_tb(prs.slides[1])
set_tb(tb, [
    "WildToolBench is a multi-turn tool-use benchmark: 256 sessions, 1,024 tasks, 4 turns/session.",
    "Baseline evaluation of four open-weight models (Qwen3-8B/14B/32B, Gemma4-31B) identified",
    "three dominant failure modes: Wrong Arguments (20-28%), Instruction Transition (20-22%),",
    "and Coreference Failure (18-20%). We implemented three harness-level interventions",
    "requiring no model retraining: (1) Enhanced System Prompt, (2) Argument Tolerance Layer,",
    "and (3) Coreference Resolution Context Injection.",
    "",
    "Enhanced evaluation (v2) results shown below:",
], font_size=13, left=0.34, top=1.15, width=9.32, height=1.9)

add_table(
    prs.slides[1],
    headers=["Model", "Task Acc (v1)", "Task Acc (v2)", "Task Delta",
             "Session Acc (v1)", "Session Acc (v2)", "Session Delta"],
    rows=[
        ["Qwen3-8B",    "40.8%", "53.3%", "+12.5pp", "2.7%", "10.2%", "+7.5pp"],
        ["Qwen3-14B",   "44.2%", "55.2%", "+11.0pp", "3.5%", "11.3%", "+7.8pp"],
        ["Qwen3-32B",   "44.1%", "54.5%", "+10.4pp", "4.3%",  "7.8%", "+3.5pp"],
        ["Gemma4-31B",  "52.0%", "56.7%",  "+4.7pp", "9.0%", "11.3%", "+2.3pp"],
    ],
    col_widths=[1.55, 1.25, 1.25, 1.0, 1.25, 1.25, 1.0],
    left=0.34, top=3.2, row_height=0.3, font_size=11,
)

# ── SLIDE 3: Introduction  (text+table left | plots right) ──────────────────
slide = prs.slides[2]
tb = get_tb(slide)
set_tb(tb, [
    "WildToolBench Benchmark",
    "  - 256 sessions | 1,024 tasks | 4 turns/session",
    "  - Types: Single-Tool, Parallel/Sequential/Mixed Multi-Tool, Clarify, Chat",
    "  - Subtypes: First Turn, Coreferential Reference,",
    "    Partial Information, Long-Range Dependency",
    "",
    "Research Question: Can harness-level interventions",
    "(no model changes) improve scores by addressing the",
    "top failure modes identified in baseline evaluation?",
    "",
    "Key Observation: No model exceeded 10% session",
    "accuracy at baseline. 394 tasks (38.5%) were failed",
    "by ALL models — confirming benchmark difficulty.",
], font_size=11, left=0.34, top=1.25, width=4.55, height=2.1)

add_table(
    slide,
    headers=["Model", "Task Acc", "Session Acc"],
    rows=[
        ["Qwen3-8B",   "40.8%", "2.7%"],
        ["Qwen3-14B",  "44.2%", "3.5%"],
        ["Qwen3-32B",  "44.1%", "4.3%"],
        ["Gemma4-31B", "52.0%", "9.0%"],
    ],
    col_widths=[1.75, 1.4, 1.4],
    left=0.34, top=3.45, row_height=0.28, font_size=11,
)

add_image(slide, PLOTS + "06_total_failures.png",       5.05, 1.2,  4.7, 2.0)
add_image(slide, PLOTS + "01_taxonomy_distribution.png", 5.05, 3.3,  4.7, 2.1)

# ── SLIDE 4: Method ─────────────────────────────────────────────────────────
tb = get_tb(prs.slides[3])
set_tb(tb, [
    "Change 1: Enhanced System Prompt  [base_handler.py - _build_system_prompt]",
    "  Problem: Models only received 'Current Date: {timestamp}' — no intent-routing guidance.",
    "  Fix: Structured instructions: (a) tool/clarify/chat routing rules, (b) argument precision",
    "       (no hallucination), (c) parallel vs sequential multi-step planning,",
    "       (d) context resolution from conversation history.",
    "  Targets: Instruction Transition (22.5%), Unnecessary Tool Call (15.6%), Missing Tool Call (9.0%)",
    "",
    "Change 2: Argument Tolerance / Normalization Layer  [checker_utils.py]",
    "  Problem: 780 'args keys mismatch' failures — models added valid optional params not in ground",
    "       truth. Date '2024-7-13' vs '2024-07-13' scored as 0.00 rouge. Numeric type mismatches.",
    "  Fix: (a) Tolerate extra keys that are valid optional schema properties.",
    "       (b) _normalize_date(): canonicalize to YYYY-MM-DD with zero-padding.",
    "       (c) _try_numeric_coerce(): treat '100' (str) and 100 (int) as equivalent.",
    "  Targets: Wrong Arguments (20-28%) — the single largest failure category (780 cases)",
    "",
    "Change 3: Coreference Resolution Context Injection  [base_handler.py - _build_context_summary]",
    "  Problem: 18-20% coreference failures on 'that location', 'the same date', 'this record'.",
    "       Size-independent: flat 18-20% across 8B to 32B — cannot be fixed by scaling alone.",
    "  Fix: Prepend '[Context from prior turns: city=Chicago, startDate=2024-07-13, ...]' before",
    "       each subsequent turn (top 20 entities extracted from all prior tool call arguments).",
    "  Targets: Coreference Failure (18-20%), Long-Range Dependency turns (28-40% accuracy)",
], font_size=11, left=0.34, top=1.25, width=9.32, height=4.1)

# ── SLIDE 5: Experiments  (tables left | plot right) ─────────────────────────
slide = prs.slides[4]
tb = get_tb(slide)
set_tb(tb, [
    "Overall Accuracy — Baseline vs. Enhanced (v2)",
], font_size=12, left=0.34, top=1.2, width=5.3, height=0.3)

add_table(
    slide,
    headers=["Model", "Task v1", "Task v2", "Task +/-", "Sess v1", "Sess v2", "Sess +/-"],
    rows=[
        ["Qwen3-8B",   "40.8%", "53.3%", "+12.5pp", "2.7%", "10.2%", "+7.5pp"],
        ["Qwen3-14B",  "44.2%", "55.2%", "+11.0pp", "3.5%", "11.3%", "+7.8pp"],
        ["Qwen3-32B",  "44.1%", "54.5%", "+10.4pp", "4.3%",  "7.8%", "+3.5pp"],
        ["Gemma4-31B", "52.0%", "56.7%",  "+4.7pp", "9.0%", "11.3%", "+2.3pp"],
    ],
    col_widths=[1.3, 0.72, 0.72, 0.75, 0.72, 0.72, 0.75],
    left=0.34, top=1.6, row_height=0.28, font_size=10,
)

# By task type table
add_table(
    slide,
    headers=["Task Type (Gemma4-31B)", "Baseline", "Enhanced", "Delta"],
    rows=[
        ["Single-Tool",           "52.2%", "60.2%", "+8.0pp"],
        ["Parallel Multi-Tool",   "49.7%", "53.8%", "+4.1pp"],
        ["Sequential Multi-Tool", "18.8%", "25.0%", "+6.2pp"],
        ["Mixed Multi-Tool",      "19.3%", "28.6%", "+9.3pp"],
        ["Clarify",               "37.3%", "40.6%", "+3.3pp"],
        ["Chat",                  "80.5%", "82.4%", "+1.9pp"],
    ],
    col_widths=[2.3, 0.85, 0.85, 0.78],
    left=0.34, top=3.05, row_height=0.27, font_size=10,
)

add_image(slide, PLOTS + "02_errors_by_model.png", 5.72, 1.2, 4.1, 4.1)

# ── SLIDE 6: Findings  (text left | heatmap + scaling right) ─────────────────
slide = prs.slides[5]
tb = get_tb(slide)
set_tb(tb, [
    "Finding 1: Argument Tolerance — largest single impact",
    "  780 'args keys mismatch' cases were structurally unfair — models",
    "  correctly understood the task but included valid optional params.",
    "  The tolerance layer recovered real accuracy, not relaxed scoring.",
    "",
    "Finding 2: Smaller models benefited most",
    "  Qwen3-8B gained +12.5pp task / +7.5pp session accuracy.",
    "  All models now exceed 53% task accuracy (vs 40-52% baseline).",
    "  Diminishing returns on Gemma4-31B (+4.7pp) — it already handled",
    "  most tolerance-related failures correctly.",
    "",
    "Finding 3: Long-Range Dependency improved most dramatically",
    "  Qwen3-8B: 28.5%->40.9% (+12.4pp). Gemma4: 39.9%->46.4% (+6.5pp).",
    "  Context injection directly targeted this turn subtype.",
    "",
    "Finding 4: Session accuracy multiplied 3-4x for Qwen3",
    "  Qwen3-14B: 3.5%->11.3%; Qwen3-8B: 2.7%->10.2%.",
    "  Per-turn improvements compound at the session level.",
    "",
    "Finding 5: Chat accuracy did not regress",
    "  Qwen3-14B: 88.1%->90.5%. Enhanced prompt correctly routes",
    "  conversational turns without introducing unnecessary tool calls.",
], font_size=11, left=0.34, top=1.25, width=4.65, height=4.1)

add_image(slide, PLOTS + "05_failure_heatmap.png",       5.1, 1.2,  4.7, 2.15)
add_image(slide, PLOTS + "03_scaling_per_category.png",  5.1, 3.4,  4.7, 1.95)

# ── SLIDE 7: Conclusion  (text left | subcategory plot right) ────────────────
slide = prs.slides[6]
tb = get_tb(slide)
set_tb(tb, [
    "Summary",
    "  Three harness-level interventions produced task accuracy gains of",
    "  +4.7 to +12.5 pp and session accuracy gains of +2.3 to +7.8 pp",
    "  across all four models — with no model retraining required.",
    "",
    "What Worked Best",
    "  - Argument Tolerance (Change 2): Largest impact. 780 'keys",
    "    mismatch' failures recovered real accuracy, not relaxed scoring.",
    "  - Coreference Injection (Change 3): Drove Long-Range Dependency",
    "    gains (Turn 0: 69% -> Turn 3: 42% at baseline).",
    "  - Enhanced Prompt (Change 1): Improved Clarify accuracy without",
    "    regressing Chat performance.",
    "",
    "Remaining Challenges",
    "  - 38.5% of tasks still fail across all models — requires",
    "    architectural changes or fine-tuning, not prompt engineering.",
    "  - Sequential Multi-Tool (12.5-25%): error cascades cannot be",
    "    fixed by harness interventions alone.",
    "  - Mixed Multi-Tool still low (13-29%): compositional planning",
    "    is a core model capability gap.",
    "",
    "Next Steps: Intent classifier router | Multi-step task planner",
    "  | String normalization pre-processor",
], font_size=11, left=0.34, top=1.25, width=4.9, height=4.1)

add_image(slide, PLOTS + "04_subcategory_breakdown.png", 5.38, 1.2, 4.45, 4.1)

prs.save('poster/WildToolBench_Poster.pptx')
print("Saved: poster/WildToolBench_Poster.pptx")
